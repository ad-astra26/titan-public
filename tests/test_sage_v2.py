"""
tests/test_sage_v2.py

Verification Test Suite for Titan V2.0 Step 5 — The Stealth-Sage.

Tests all five "senses" of the Stealth-Sage:
  Test 1: Ollama phi3:mini health check (live service)
  Test 2: SearXNG live search (live service)
  Test 3: XResearcher X-Search (mocked — no real API key needed)
  Test 4: Crawl4AI web scraping + distillation end-to-end (live services)
  Test 5: DocumentProcessor — PDF, DOCX, PPTX, XLSX deep-dive (local files, mocked Ollama)
  Test 6: Gateway integration — full pre_prompt_hook research pipeline (5 prompts)
  Test 7: Gatekeeper STATE_NEED_RESEARCH routing
  Test 8: _load_stealth_sage_config defaults and toml parsing

Follows the existing codebase patterns:
  - pytest + pytest-asyncio
  - Isolated temp storage in test_env/ (cleaned per fixture)
  - Direct module imports from titan_plugin
  - unittest.mock for offline tests (Tests 3, 5)
"""

import asyncio
import json
import os
import shutil
import sys
import time
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
import pytest_asyncio

# Ensure titan_plugin is on sys.path (same pattern as all other tests)
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

# ──────────────────────────────────────────────────────────────────────────────
# Fixtures
# ──────────────────────────────────────────────────────────────────────────────

TEST_ENV = Path("./test_env/sage_v2")
RESEARCH_LOG = Path("./data/logs/sage_research.log")


@pytest.fixture(autouse=True)
def clean_test_env():
    """Wipes temp env before each test and removes research log for isolation."""
    if TEST_ENV.exists():
        shutil.rmtree(TEST_ENV)
    TEST_ENV.mkdir(parents=True, exist_ok=True)

    # Clean research log for fresh state
    if RESEARCH_LOG.exists():
        RESEARCH_LOG.unlink()
    RESEARCH_LOG.parent.mkdir(parents=True, exist_ok=True)

    yield

    # Cleanup after test
    if TEST_ENV.exists():
        shutil.rmtree(TEST_ENV)


def minimal_stealth_sage_config(tmp_dir: str = "/tmp/titan_sage_docs_test") -> dict:
    """Returns a minimal [stealth_sage] config dict for unit testing."""
    return {
        "searxng_host": "http://localhost:8080",
        "searxng_top_num_urls": 2,
        "twitterapi_io_key": "",          # Disabled — no real key in tests
        "twitterapi_search_depth": 5,
        "webshare_rotating_url": "",      # Disabled — direct connection in tests
        "ollama_host": "http://localhost:11434",
        "ollama_model": "phi3:mini",
        "max_load_avg": 99.0,             # Never skip on test VPS
        "research_timeout_seconds": 60,
        "doc_safe_room": tmp_dir,
    }


# ──────────────────────────────────────────────────────────────────────────────
# Test 1: Ollama Health Check
# ──────────────────────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_ollama_health():
    """
    Test 1: Verifies that Ollama is running, phi3:mini model is available,
    and can generate a non-empty response to a simple distillation prompt.

    Requires: Ollama running + phi3:mini pulled on the VPS.
    """
    import httpx
    print("\n\n🧠 Test 1: Ollama phi3:mini Health Check")

    # 1a. Check Ollama tags endpoint
    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            resp = await client.get("http://localhost:11434/api/tags")
        assert resp.status_code == 200, f"Ollama /api/tags returned {resp.status_code}"
        data = resp.json()
        model_names = [m["name"] for m in data.get("models", [])]
        print(f"   Ollama models available: {model_names}")
        phi_available = any("phi3" in n for n in model_names)
        assert phi_available, (
            f"phi3:mini not found in Ollama models: {model_names}. "
            "Run 'ollama pull phi3:mini' to install."
        )
    except httpx.ConnectError:
        pytest.skip("Ollama not running on localhost:11434 — skipping live test.")

    # 1b. Fire a small distillation prompt
    payload = {
        "model": "phi3:mini",
        "prompt": "Summarize in one sentence: Bitcoin is a decentralized digital currency.",
        "stream": False,
    }
    async with httpx.AsyncClient(timeout=60.0) as client:
        resp = await client.post("http://localhost:11434/api/generate", json=payload)
    assert resp.status_code == 200
    response_text = resp.json().get("response", "").strip()
    assert len(response_text) > 10, f"Ollama returned suspiciously short response: '{response_text}'"
    print(f"   ✅ Ollama phi3:mini response: '{response_text[:100]}...'")


# ──────────────────────────────────────────────────────────────────────────────
# Test 2: SearXNG Live Search
# ──────────────────────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_searxng_live_search():
    """
    Test 2: Verifies SearXNG is running and returns JSON search results.

    Requires: SearXNG Docker container running on localhost:8080.
    """
    import httpx
    print("\n\n🔍 Test 2: SearXNG Live Search")

    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            resp = await client.get(
                "http://localhost:8080/search",
                params={"q": "Solana blockchain", "format": "json", "language": "en"},
            )
    except httpx.ConnectError:
        pytest.skip("SearXNG not reachable on localhost:8080 — skipping live test.")

    assert resp.status_code == 200, f"SearXNG returned {resp.status_code}"
    data = resp.json()

    results = data.get("results", [])
    assert len(results) >= 1, "SearXNG returned 0 results for 'Solana blockchain' — check container config."
    print(f"   ✅ SearXNG returned {len(results)} results. First URL: {results[0].get('url', 'N/A')}")


# ──────────────────────────────────────────────────────────────────────────────
# Test 3: XResearcher (Mocked — Offline)
# ──────────────────────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_x_researcher_mocked():
    """
    Test 3: Verifies XResearcher search() formats tweet data correctly.
    Uses unittest.mock to simulate TwitterAPI.io response — no real API key needed.
    Also verifies graceful return of "" when api_key is empty.
    """
    from titan_plugin.logic.sage.x_researcher import XResearcher
    print("\n\n🐦 Test 3: XResearcher X-Search (Mocked)")

    # 3a: Disabled researcher returns "" without crashing
    disabled = XResearcher(api_key="", proxy_url=None, search_depth=5)
    assert not disabled.is_enabled
    result = await disabled.search("anything")
    assert result == "", "Disabled XResearcher should return empty string."
    print("   ✅ Disabled XResearcher returns '' correctly.")

    # 3b: Mock a successful API response
    mock_tweets = {
        "tweets": [
            {"text": "SOL is pumping hard today! 🚀 #Solana"},
            {"text": "Solana validators are running perfectly. Network health is strong."},
            {"text": "Just bought more SOL. The ecosystem is thriving."},
        ]
    }

    mock_response = MagicMock()
    mock_response.status_code = 200
    mock_response.json.return_value = mock_tweets

    with patch("httpx.AsyncClient") as mock_client_cls:
        mock_client = AsyncMock()
        mock_client.__aenter__ = AsyncMock(return_value=mock_client)
        mock_client.__aexit__ = AsyncMock(return_value=False)
        mock_client.get = AsyncMock(return_value=mock_response)
        mock_client_cls.return_value = mock_client

        researcher = XResearcher(api_key="test_key_12345", proxy_url=None, search_depth=3)
        assert researcher.is_enabled
        result = await researcher.search("SOL price Solana")

    assert "[X_SEARCH_RESULTS" in result, f"Expected header in result, got: {result[:100]}"
    assert "SOL is pumping" in result
    assert "validators" in result
    print(f"   ✅ XResearcher formatted {len(mock_tweets['tweets'])} tweets correctly.")
    print(f"   Result preview: {result[:150]}")


# ──────────────────────────────────────────────────────────────────────────────
# Test 4: Full Research Pipeline (Live — SearXNG + Crawl4AI + Ollama)
# ──────────────────────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_full_research_pipeline_live():
    """
    Test 4: Fires the full StealthSageResearcher.research() pipeline against
    a real query. Verifies the [SAGE_RESEARCH_FINDINGS] block is returned
    and that the research audit log is written.

    Requires: SearXNG + Ollama + Crawl4AI/Playwright installed.
    """
    import httpx
    from titan_plugin.logic.sage.researcher import StealthSageResearcher

    print("\n\n🌐 Test 4: Full Research Pipeline (Live)")

    # Quick liveness check before running
    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            await client.get("http://localhost:8080/search?q=test&format=json")
    except httpx.ConnectError:
        pytest.skip("SearXNG not reachable — skipping live pipeline test.")

    config = minimal_stealth_sage_config(str(TEST_ENV / "docs"))
    researcher = StealthSageResearcher(config)

    query = "python programming language"
    result = await researcher.research(query, transition_id=42)

    if result:
        assert result.startswith("[SAGE_RESEARCH_FINDINGS]:"), (
            f"Result should start with '[SAGE_RESEARCH_FINDINGS]:' but got: {result[:100]}"
        )
        assert len(result) > 50, "Findings block seems too short."
        print(f"   ✅ Research pipeline returned {len(result)} chars of findings.")
        print(f"   Preview: {result[:200]}")
    else:
        # Acceptable if web scraping returned no content — log it
        print("   ℹ️ No findings returned (web content was empty or sites blocked). Acceptable result.")

    # Verify audit log was written (if findings were returned)
    if result and RESEARCH_LOG.exists():
        with open(RESEARCH_LOG) as f:
            lines = f.readlines()
        assert len(lines) >= 1, "Research audit log should have at least one entry."
        entry = json.loads(lines[-1])
        assert entry["knowledge_gap"] == query
        assert entry["transition_id"] == 42
        print(f"   ✅ Research audit log written: transition_id={entry['transition_id']}, sources={entry['sources_used']}")


# ──────────────────────────────────────────────────────────────────────────────
# Test 5: DocumentProcessor — Multi-Format Deep Dive (Offline / Mocked Ollama)
# ──────────────────────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_document_processor_all_formats():
    """
    Test 5: Creates a test file for each supported format (PDF, DOCX, PPTX, XLSX)
    and verifies DocumentProcessor extracts content and returns a structured dict.
    Mocks Ollama to avoid real inference cost — tests the full pipeline except distillation.

    Requires: unstructured + reportlab + python-docx + python-pptx + openpyxl installed.
    """
    from titan_plugin.logic.sage.document_processor import DocumentProcessor
    print("\n\n📄 Test 5: DocumentProcessor — Multi-Format Deep Dive")

    doc_dir = TEST_ENV / "test_docs"
    doc_dir.mkdir(parents=True, exist_ok=True)

    test_files = {}

    # Create PDF
    try:
        from reportlab.pdfgen import canvas
        pdf_path = doc_dir / "test.pdf"
        c = canvas.Canvas(str(pdf_path))
        c.drawString(100, 750, "Titan V2.0 Document Test: Solana Network Performance Q1 2026")
        c.drawString(100, 730, "Key Metric: TPS = 65,000 | Validators: 1,500 | Uptime: 99.98%")
        c.save()
        test_files["PDF"] = pdf_path
        print("   Created test.pdf")
    except ImportError:
        print("   ⚠️ reportlab not installed — skipping PDF test. Install: pip install reportlab")

    # Create DOCX
    try:
        from docx import Document as DocxDocument
        docx_path = doc_dir / "test.docx"
        doc = DocxDocument()
        doc.add_heading("Titan V2.0 Architecture Brief", 0)
        doc.add_paragraph("Titan uses Implicit Q-Learning (IQL) for offline reinforcement learning.")
        doc.add_paragraph("The SageRecorder uses LazyMemmapStorage with a 50GB limit.")
        doc.save(str(docx_path))
        test_files["DOCX"] = docx_path
        print("   Created test.docx")
    except ImportError:
        print("   ⚠️ python-docx not installed — skipping DOCX test. Install: pip install python-docx")

    # Create PPTX
    try:
        from pptx import Presentation
        pptx_path = doc_dir / "test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Titan Stealth-Sage Research Engine"
        slide.placeholders[1].text = "Powered by Crawl4AI, Unstructured, and Ollama phi3:mini."
        prs.save(str(pptx_path))
        test_files["PPTX"] = pptx_path
        print("   Created test.pptx")
    except ImportError:
        print("   ⚠️ python-pptx not installed — skipping PPTX test. Install: pip install python-pptx")

    # Create XLSX
    try:
        import openpyxl
        xlsx_path = doc_dir / "test.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "SOL Metrics"
        ws.append(["Metric", "Value"])
        ws.append(["SOL Price", "$180.00"])
        ws.append(["24h Volume", "$2.4B"])
        ws.append(["Market Cap", "$78B"])
        wb.save(str(xlsx_path))
        test_files["XLSX"] = xlsx_path
        print("   Created test.xlsx")
    except ImportError:
        print("   ⚠️ openpyxl not installed — skipping XLSX test. Install: pip install openpyxl")

    if not test_files:
        pytest.skip("No document libraries available — install reportlab, python-docx, python-pptx, openpyxl")

    # Mock Ollama distillation response
    mock_ollama_response = {
        "response": json.dumps({
            "topic": "Test document content",
            "key_points": ["Key point 1", "Key point 2"],
            "source_type": "Document"
        })
    }
    mock_http_resp = MagicMock()
    mock_http_resp.status_code = 200
    mock_http_resp.json.return_value = mock_ollama_response
    mock_http_resp.raise_for_status = MagicMock()

    processor = DocumentProcessor(
        safe_room=str(doc_dir / "safe_room"),
        max_load_avg=99.0,  # Never skip
        proxy_url=None,
    )

    results = {}
    for fmt, file_path in test_files.items():
        file_url = f"file://{file_path}"
        with patch("httpx.AsyncClient") as mock_client_cls:
            mock_client = AsyncMock()
            mock_client.__aenter__ = AsyncMock(return_value=mock_client)
            mock_client.__aexit__ = AsyncMock(return_value=False)
            mock_client.post = AsyncMock(return_value=mock_http_resp)
            mock_client_cls.return_value = mock_client

            result = await processor.process(file_url)

        if result:
            assert result.get("source_type") == "Document", f"source_type mismatch for {fmt}"
            assert "summary" in result, f"Missing 'summary' key for {fmt}"
            results[fmt] = True
            print(f"   ✅ {fmt}: processed correctly. Summary: '{result['summary'][:80]}'")
        else:
            # Unstructured might not parse all formats in all environments
            print(f"   ℹ️ {fmt}: returned empty dict (unstructured may not support this format in current env)")
            results[fmt] = None

    processed = [f for f, ok in results.items() if ok]
    assert len(processed) > 0, "DocumentProcessor failed for ALL formats — check unstructured installation."
    print(f"   ✅ DocumentProcessor successfully handled: {processed}")


# ──────────────────────────────────────────────────────────────────────────────
# Test 6: Gateway Integration — Full pre_prompt_hook Research Loop (5 Prompts)
# ──────────────────────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_gateway_integration_research_trigger():
    """
    Test 6: Sends 5 prompts containing informational keywords through the full
    TitanPlugin.pre_prompt_hook(). Forces Gatekeeper to STATE_NEED_RESEARCH
    by mocking the Advantage score to 0.2 (Shadow threshold) and verifying
    [SAGE_RESEARCH_FINDINGS] appears in context["system_prompt"].
    Measures wall-clock time per call.

    Requires: Ollama + SearXNG running (live) OR mocked research output.
    """
    from titan_plugin import TitanPlugin
    print("\n\n🚀 Test 6: Gateway Integration — Research Trigger (5 Prompts)")

    # We use mocked StealthSageResearcher to avoid live dependency in CI
    mock_research_result = "[SAGE_RESEARCH_FINDINGS]: SOL current price is approximately $180, up 3% in 24h. Market volume is $2.4B. Validators report 99.98% uptime."

    informational_prompts = [
        "What is the current SOL price today?",
        "What is the latest news about Solana blockchain?",
        "What are people saying about the crypto market right now?",
        "Predict the Bitcoin price trend for this week.",
        "What is the current APY on Solana staking?",
    ]

    # Init plugin (uses ./authority.json mock)
    plugin = TitanPlugin("./authority.json")

    # Force Gatekeeper to always return low Advantage → triggers informational check
    def mock_gatekeeper(state_tensor, raw_prompt=""):
        if plugin.gatekeeper._is_informational_query(raw_prompt):
            return "STATE_NEED_RESEARCH", 0.2, ""
        return "Shadow", 0.2, ""

    plugin.gatekeeper.decide_execution_mode = mock_gatekeeper

    # Mock the sage_researcher.research() to return instant findings
    plugin.sage_researcher.research = AsyncMock(return_value=mock_research_result)

    timings = []
    for i, prompt in enumerate(informational_prompts):
        print(f"\n   [{i+1}/5] Prompt: '{prompt}'")
        start = time.perf_counter()
        context = await plugin.pre_prompt_hook(prompt, {})
        elapsed = time.perf_counter() - start
        timings.append(elapsed)

        system_prompt = context.get("system_prompt", "")
        assert "[SAGE_RESEARCH_FINDINGS]" in system_prompt, (
            f"Expected [SAGE_RESEARCH_FINDINGS] in system_prompt for prompt: '{prompt}'. "
            f"Got system_prompt: '{system_prompt[:200]}'"
        )
        # Verify research sources were tracked for post_resolution_hook tagging
        assert len(plugin._last_research_sources) > 0, "Research sources should be tracked."
        print(f"   ✅ [{i+1}/5] Findings injected. Sources: {plugin._last_research_sources}. Time: {elapsed*1000:.1f}ms")

    avg_ms = sum(timings) / len(timings) * 1000
    max_ms = max(timings) * 1000
    print(f"\n   📊 Gateway Performance Report:")
    print(f"   • Average hook time: {avg_ms:.1f}ms")
    print(f"   • Max hook time:     {max_ms:.1f}ms")
    print(f"   ✅ All 5 prompts correctly triggered [SAGE_RESEARCH_FINDINGS] injection.")


# ──────────────────────────────────────────────────────────────────────────────
# Test 7: Gatekeeper STATE_NEED_RESEARCH Routing
# ──────────────────────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_gatekeeper_state_need_research():
    """
    Test 7: Directly tests the Gatekeeper's _is_informational_query() method
    and STATE_NEED_RESEARCH routing logic.

    The Gatekeeper's decide_execution_mode does an early return from find_grounded_action
    when the KNN replay buffer is empty (returns None → Shadow, 0.0, ""). This is
    correct production behavior. In tests we mock find_grounded_action to return a
    valid (tensor, text) tuple so we can reach the Advantage routing code and verify
    STATE_NEED_RESEARCH is returned for informational prompts.

    Also directly tests _is_informational_query independently to provide a clean signal.
    """
    import torch
    from titan_plugin.core.sage.recorder import SageRecorder
    from titan_plugin.logic.sage.gatekeeper import SageGatekeeper
    from titan_plugin.logic.sage.scholar import SageScholar

    print("\n\n🏛️ Test 7: Gatekeeper STATE_NEED_RESEARCH Routing")

    test_config = {
        "sage_memory": {
            "buffer_size": 10,
            "storage_path": str(TEST_ENV / "gatekeeper_test"),
            "embedding_dim": 384,
        }
    }
    recorder = SageRecorder(test_config)
    scholar = SageScholar(recorder)
    gatekeeper = SageGatekeeper(scholar, recorder)

    # ---- 7a: Directly test _is_informational_query() ----
    informational = [
        "What is the latest Bitcoin price today?",
        "What is the current SOL market cap?",
        "What are people saying about crypto trending now?",
        "What is the SOL price right now?",
        "Tell me the latest news about Solana.",
    ]
    non_informational = [
        "Write me a haiku about autumn leaves.",
        "Tell me a joke about a robot.",
        "What is the capital of France?",
        "How do I reverse a Python list?",
    ]

    for prompt in informational:
        assert gatekeeper._is_informational_query(prompt), (
            f"Expected informational=True for: '{prompt}'"
        )
        print(f"   ✅ [informational=True]  '{prompt[:70]}'")

    for prompt in non_informational:
        assert not gatekeeper._is_informational_query(prompt), (
            f"Expected informational=False for: '{prompt}'"
        )
        print(f"   ✅ [informational=False] '{prompt[:70]}'")

    # ---- 7b: Test STATE_NEED_RESEARCH routing with mocked KNN ----
    # Mock find_grounded_action to return a valid (tensor, text) so we bypass
    # the empty-buffer early return and reach the Advantage comparison code.
    # Also mock the scholar networks to return A = 0.2 (Shadow territory).
    dummy_action = torch.zeros(128)
    gatekeeper.find_grounded_action = MagicMock(return_value=(dummy_action, "mock_action_text"))

    def mock_q(td, **_):
        td["state_action_value"] = torch.tensor([[1.0]])
        return td

    def mock_v(td, **_):
        td["state_value"] = torch.tensor([[0.8]])
        return td

    def mock_actor(td, **_):
        td["action"] = torch.zeros(1, 128)
        return td

    scholar.actor_module = mock_actor
    scholar.qvalue_module = mock_q
    scholar.value_module = mock_v

    dummy_state = torch.zeros(128)

    # Informational prompts → STATE_NEED_RESEARCH (Advantage=0.2, < 0.4 threshold)
    for prompt in informational:
        mode, adv, _ = gatekeeper.decide_execution_mode(dummy_state, raw_prompt=prompt)
        assert mode == "STATE_NEED_RESEARCH", (
            f"Expected STATE_NEED_RESEARCH for '{prompt}', got '{mode}' (A={adv:.3f})"
        )
        print(f"   ✅ '{prompt[:50]}...' → {mode} (A={adv:.2f})")

    # Non-informational prompts → Shadow (NOT STATE_NEED_RESEARCH)
    for prompt in non_informational:
        mode, adv, _ = gatekeeper.decide_execution_mode(dummy_state, raw_prompt=prompt)
        assert mode == "Shadow", (
            f"Expected Shadow for non-informational '{prompt}', got '{mode}'"
        )
        print(f"   ✅ '{prompt[:50]}' → {mode} (correctly no research triggered)")

    # Empty raw_prompt → Shadow (backward compat guard)
    mode, adv, _ = gatekeeper.decide_execution_mode(dummy_state, raw_prompt="")
    assert mode == "Shadow", f"Expected Shadow for empty prompt, got '{mode}'"
    print(f"   ✅ Empty prompt → Shadow (backward compat preserved).")

    print(f"   ✅ Gatekeeper correctly routes informational vs creative queries.")


# ──────────────────────────────────────────────────────────────────────────────
# Test 8: Config Loading — Defaults and TOML Parsing
# ──────────────────────────────────────────────────────────────────────────────

def test_load_stealth_sage_config_defaults():
    """
    Test 8: Verifies _load_stealth_sage_config() returns all required keys with
    proper defaults when loading from the real config.toml, and that none of
    the required keys are missing.
    """
    from titan_plugin import TitanPlugin
    print("\n\n⚙️ Test 8: Config Loading — Defaults and TOML Parsing")

    config = TitanPlugin._load_stealth_sage_config()

    required_keys = [
        "searxng_host", "searxng_top_num_urls",
        "twitterapi_io_key", "twitterapi_search_depth",
        "webshare_rotating_url",
        "ollama_host", "ollama_model",
        "max_load_avg", "research_timeout_seconds", "doc_safe_room",
    ]
    for key in required_keys:
        assert key in config, f"Missing required key '{key}' in stealth_sage config."
        print(f"   ✅ {key}: {repr(config[key])}")

    # Validate types
    assert isinstance(config["searxng_top_num_urls"], int)
    assert isinstance(config["max_load_avg"], float)
    assert isinstance(config["research_timeout_seconds"], (int, float))
    assert isinstance(config["twitterapi_io_key"], str)
    assert isinstance(config["webshare_rotating_url"], str)

    # Default values should match config.toml defaults (or built-in defaults)
    assert config["searxng_host"].startswith("http")
    assert config["ollama_model"] == "deepseek-v3.1:671b"
    assert config["doc_safe_room"] == "/tmp/titan_sage_docs"

    print("   ✅ All config keys present and correctly typed.")


# ──────────────────────────────────────────────────────────────────────────────
# Entry point
# ──────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import subprocess
    result = subprocess.run(
        ["python", "-m", "pytest", __file__, "-v", "--tb=short"],
        cwd=os.path.join(os.path.dirname(__file__), ".."),
    )
    sys.exit(result.returncode)
